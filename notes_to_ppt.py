import re
import json
from pptx import Presentation
import requests
from json import JSONDecodeError

prompt = """
Turn these notes into a detailed PowerPoint outline with multiple bullet points per slide.
Output ONLY valid JSON as an array of slides.
Each slide must have: "title" and "bullets".
Each bullet must have "text" and optional "description".

Create 3-5 bullet points for each slide with detailed explanations.

Example format:
[
  {
    "title": "Slide Title",
    "bullets": [
      {"text": "Main point 1", "description": "Detailed explanation 1"},
      {"text": "Main point 2", "description": "Detailed explanation 2"},
      {"text": "Main point 3", "description": "Detailed explanation 3"}
    ]
  }
]

Notes to expand into detailed slides:
- God answers prayer: Explain how prayer works, types of prayers, examples, biblical references
- Faith is necessary: Define faith, importance of faith, how to develop faith, biblical examples
- Patience brings results: Benefits of patience, how patience leads to results, biblical stories about patience
"""

# --- Call Ollama API ---
url = "http://localhost:11434/api/generate"
payload = {"model": "llama2", "prompt": prompt, "stream": False}
response = requests.post(url, json=payload)
response.raise_for_status()
output = response.json()["response"]

print("=== RAW OUTPUT ===")
print(repr(output))
print("==================")

# --- Clean output ---
# Try to extract JSON content
json_match = re.search(r'(\[.*\]|\{.*\})', output, re.DOTALL)
if json_match:
    clean_output = json_match.group(1)
else:
    # Fallback: try basic cleaning
    output = re.sub(r"```.*?```", "", output, flags=re.DOTALL)
    output = re.sub(r",(\s*[\}\]])", r"\1", output)
    output_lines = [line for line in output.splitlines() if not line.strip().lower().startswith("note:")]
    clean_output = "\n".join(output_lines).strip()

# Remove trailing commas
clean_output = re.sub(r',\s*}', '}', clean_output)
clean_output = re.sub(r',\s*]', ']', clean_output)

print("=== CLEANED OUTPUT ===")
print(repr(clean_output))
print("======================")

# --- Parse JSON ---
try:
    slides_json = json.loads(clean_output)
except JSONDecodeError as e:
    print(f"JSON parsing failed: {e}")
    # Try one more cleanup
    clean_output = re.sub(r'^[^{[]*', '', clean_output)  # Remove text before { or [
    clean_output = re.sub(r'[^}]*$', '', clean_output)   # Remove text after } or ]
    clean_output = clean_output.strip()
    print(f"Retrying with: {clean_output}")
    slides_json = json.loads(clean_output)

# Handle different JSON structures
if isinstance(slides_json, list):
    slides_data = slides_json
elif isinstance(slides_json, dict) and "slides" in slides_json:
    slides_data = slides_json["slides"]
else:
    slides_data = []

# --- Build PowerPoint ---
prs = Presentation()
for slide in slides_data:
    title = slide.get("title", "Untitled Slide")
    bullets = slide.get("bullets", [])
    slide_layout = prs.slide_layouts[1]  # Title + Content
    s = prs.slides.add_slide(slide_layout)
    s.shapes.title.text = title
    body = s.placeholders[1].text_frame
    
    # Clear any existing paragraphs first
    body.clear()
    
    # Add bullet points
    for i, bullet in enumerate(bullets):
        text = bullet.get("text", "")
        description = bullet.get("description", "")
        
        # Add main bullet point
        p = body.add_paragraph()
        p.text = text
        p.level = 0
        
        # Add description as sub-bullet if it exists
        if description:
            sub_p = body.add_paragraph()
            sub_p.text = description
            sub_p.level = 1

prs.save("notes_presentation_json.pptx")
print("✅ PowerPoint saved as notes_presentation_json.pptx")
print(f"Created {len(slides_data)} slides with total of {sum(len(slide.get('bullets', [])) for slide in slides_data)} bullet points")